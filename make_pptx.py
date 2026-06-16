"""Generates Becca-Pitch.pptx — a branded, animated overview deck."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

HERE = os.path.dirname(__file__)
LOGO = os.path.join(HERE, "assets", "logo.png")
BANNER = os.path.join(HERE, "assets", "banner.png")
OUT = os.path.join(HERE, "..", "..", "Becca-Pitch.pptx")

PINK = RGBColor(0xD6, 0x33, 0x6C)
BLUE = RGBColor(0x1C, 0x7E, 0xD6)
PLUM = RGBColor(0x8A, 0x5A, 0xC8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0x6B, 0x5B, 0x69)
INK = RGBColor(0x3A, 0x2A, 0x38)
SOFT = RGBColor(0xFF, 0xF4, 0xF9)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def grad_bg(slide, c1, c2, angle=45):
    fill = slide.background.fill
    fill.gradient()
    try:
        fill.gradient_angle = angle
    except Exception:
        pass
    s = fill.gradient_stops
    s[0].position = 0.0
    s[0].color.rgb = RGBColor(*c1)
    s[1].position = 1.0
    s[1].color.rgb = RGBColor(*c2)


def solid_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, l, t, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, (text, size, color, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(6)
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = "Arial"
    return tb


TRANS = {
    "fade": ("fade", {}),
    "pushL": ("push", {"dir": "l"}),
    "coverD": ("cover", {"dir": "d"}),
    "wipeR": ("wipe", {"dir": "r"}),
    "dissolve": ("dissolve", {}),
    "newsflash": ("newsflash", {}),
    "circle": ("circle", {}),
}


def transition(slide, key="fade", speed="med"):
    name, attrs = TRANS[key]
    tr = etree.SubElement(slide._element, qn("p:transition"))
    tr.set("spd", speed)
    child = etree.SubElement(tr, qn(f"p:{name}"))
    for k, v in attrs.items():
        child.set(k, v)


def slide(bg_fn):
    s = prs.slides.add_slide(BLANK)
    bg_fn(s)
    return s


# 1) COVER
s = slide(lambda x: grad_bg(x, (0xFF, 0x8F, 0xB1), (0x7F, 0xD4, 0xFF), 45))
s.shapes.add_picture(LOGO, Inches((13.333 - 2.6) / 2), Inches(1.1), height=Inches(2.6))
box(s, 1, 3.9, 11.333, 2.6, [
    ("Becca", 66, WHITE, True),
    ("Less paperwork. More playtime.", 28, WHITE, False),
    ("A WhatsApp helper that writes your ECE documentation for you.", 18, WHITE, False),
], align=PP_ALIGN.CENTER)
transition(s, "fade")

# 2) THE PROBLEM
s = slide(lambda x: solid_bg(x, SOFT))
box(s, 0.9, 0.7, 11.5, 1.2, [("The problem 😮‍💨", 40, PINK, True)])
box(s, 0.9, 2.1, 11.5, 4.5, [
    ("Every term, an ECE student hand-writes a huge submission.", 24, INK, False),
    ("Term 3 alone is ~94 pages: 24 observations, 50 planning sheets,", 22, GREY, False),
    ("4 PPP sheets, a documentation board, and 3 essays.", 22, GREY, False),
    ("Same reflections, over and over, in a strict format. Exhausting.", 24, INK, True),
])
transition(s, "pushL")

# 3) HOW IT WORKS
s = slide(lambda x: solid_bg(x, WHITE))
box(s, 0.9, 0.7, 11.5, 1.2, [("How it works — 4 little steps", 40, BLUE, True)])
box(s, 0.9, 2.2, 11.5, 4.5, [
    ("1.  📸  Snap a photo of the children at play", 26, INK, True),
    ("2.  🗣️  Add one short note (names, week, big idea)", 26, INK, True),
    ("3.  ✨  Becca writes the full sheet, in proper words", 26, INK, True),
    ("4.  📄  You review, edit, and download the .docx", 26, INK, True),
])
transition(s, "coverD")

# 4) WHAT IT MAKES
s = slide(lambda x: solid_bg(x, SOFT))
box(s, 0.9, 0.7, 11.5, 1.2, [("It writes every part of the assignment 💕", 38, PINK, True)])
box(s, 0.9, 2.1, 11.5, 4.6, [
    ("📝  Observation sheets (AOR)", 24, INK, True),
    ("🧩  Planning sheets (FCPS)", 24, INK, True),
    ("🎨  PPP + observation", 24, INK, True),
    ("🖼️  Documentation board", 24, INK, True),
    ("💭  The 3 reflective essays", 24, INK, True),
])
transition(s, "wipeR")

# 5) ON WHATSAPP
s = slide(lambda x: grad_bg(x, (0x7F, 0xD4, 0xFF), (0xB1, 0x97, 0xFC), 30))
box(s, 0.9, 0.8, 11.5, 1.2, [("Just WhatsApp — no app to install 📱", 38, WHITE, True)])
box(s, 0.9, 2.3, 11.5, 4.2, [
    ("She texts a photo + note  →  it comes back as a finished file.", 26, WHITE, True),
    ("Send a PHOTO for an observation, or type:", 22, WHITE, False),
    ("plan · ppp · board · essay1 / essay2 / essay3", 24, WHITE, True),
    ("Usually under a minute. Each photo can become its own sheet.", 20, WHITE, False),
])
transition(s, "dissolve")

# 6) THE DOCUMENT
s = slide(lambda x: solid_bg(x, WHITE))
box(s, 0.9, 0.7, 11.5, 1.1, [("A polished, real document", 40, BLUE, True)])
s.shapes.add_picture(BANNER, Inches(0.9), Inches(1.9), width=Inches(7.5))
box(s, 0.9, 3.6, 11.5, 3.2, [
    ("• Becca gradient banner + framed page border", 22, INK, False),
    ("• The child's real photo embedded in the sheet", 22, INK, False),
    ("• Written like a real Canadian student — human, not AI", 22, INK, False),
    ("• Times New Roman 12pt double-spaced (submission-safe)", 22, INK, False),
])
transition(s, "circle")

# 7) ROADMAP
s = slide(lambda x: solid_bg(x, SOFT))
box(s, 0.9, 0.7, 11.5, 1.2, [("The plan 🚀", 40, PINK, True)])
box(s, 0.9, 2.1, 11.5, 4.6, [
    ("✅  Phase 1 — the magic moment (DONE & proven)", 24, INK, True),
    ("✅  Phase 2 — every document type (BUILT)", 24, INK, True),
    ("⏳  Compile the whole term into one file + progress", 24, GREY, False),
    ("⏳  24/7 hosting · real “Becca” WhatsApp number + logo", 24, GREY, False),
    ("⏳  Accounts & billing for many students", 24, GREY, False),
])
transition(s, "newsflash")

# 8) CLOSING
s = slide(lambda x: grad_bg(x, (0xFF, 0x8F, 0xB1), (0x7F, 0xD4, 0xFF), 60))
s.shapes.add_picture(LOGO, Inches((13.333 - 2.0) / 2), Inches(1.4), height=Inches(2.0))
box(s, 1, 3.8, 11.333, 2.6, [
    ("Less paperwork. More playtime. 🍁", 40, WHITE, True),
    ("Made with care for Rebecca & her little learners.", 22, WHITE, False),
], align=PP_ALIGN.CENTER)
transition(s, "fade")

prs.save(OUT)
print("saved:", os.path.abspath(OUT), "·", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
